import os
import pandas as pd
import matplotlib.pyplot as plt
import streamlit as st
from openai import OpenAI
from PIL import Image
import PyPDF2
from io import BytesIO
import docx
import pptx
import struct
import fcsparser
import numpy as np
import tempfile

st.set_page_config(page_title="File Analysis LLM Agent", layout="wide")

# ─── OpenAI Configuration ────────────────────────────────
MODEL_OPTIONS = {
    "GPT-3.5 Turbo (Fastest/Cheapest)": "gpt-3.5-turbo-0125",
    "GPT-4 Turbo (Balanced)": "gpt-4-turbo-preview",
    "GPT-4 Omni (Most Powerful)": "gpt-4o"

}
import os
api_key = os.getenv("OPENAI_API_KEY")  # ✅ Safer


with st.sidebar:
    st.subheader("API Configuration")
    api_key = st.text_input("OpenAI API Key", value=api_key, type="password")
    selected_model = st.selectbox("Select Model", list(MODEL_OPTIONS.keys()), index=0)
    st.markdown("### Cost Estimates (per 1M tokens)")
    st.markdown("- GPT-3.5: $0.50 / $1.50")
    st.markdown("- GPT-4 Turbo: $10 / $30")
    st.markdown("- GPT-4 Omni: $5 / $15")

if not api_key:
    st.error("🔑 Please enter your OpenAI API key")
    st.stop()

try:
    client = OpenAI(api_key=api_key)
    client.models.list()
    st.sidebar.success("✅ API Key Valid")
except Exception as e:
    st.sidebar.error(f"❌ API Error: {str(e)}")
    st.stop()

if "history" not in st.session_state:
    st.session_state.history = []
if "uploaded_files" not in st.session_state:
    st.session_state.uploaded_files = {}

# ─── File Parsers ─────────────────────────────────────────
def extract_text_from_pdf(file):
    reader = PyPDF2.PdfReader(file)
    return "\n".join([page.extract_text() for page in reader.pages[:10]])

def extract_text_from_docx(file):
    doc = docx.Document(file)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pptx(file):
    prs = pptx.Presentation(file)
    return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])

def parse_fcs_file(file):
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".fcs") as tmp_file:
            tmp_file.write(file.read())
            temp_path = tmp_file.name
        meta, data = fcsparser.parse(temp_path, reformat_meta=True)
        data.columns = [col.strip() for col in data.columns]
        data = data.apply(pd.to_numeric, errors='ignore')
        os.remove(temp_path)
        return data
    except Exception as e:
        return f"Error parsing FCS file: {str(e)}"

def parse_hgb_file(file):
    try:
        data = file.read()
        total_floats = len(data) // 4
        if total_floats % 3 != 0:
            return f"Invalid HGB format: total float count ({total_floats}) not divisible by 3"
        values = struct.unpack(f'{total_floats}f', data)
        chunk = total_floats // 3
        return pd.DataFrame({
            "US1": values[:chunk],
            "US2": values[chunk:2*chunk],
            "US3": values[2*chunk:]
        })
    except Exception as e:
        return f"Error parsing HGB file: {str(e)}"

# ─── File Handling ───────────────────────────────────────
def process_uploaded_files(uploaded_files):
    for file in uploaded_files:
        if file.name not in st.session_state.uploaded_files:
            ext = file.name.split('.')[-1].lower()
            try:
                if ext == 'pdf':
                    content = extract_text_from_pdf(BytesIO(file.getvalue()))
                elif ext == 'docx':
                    content = extract_text_from_docx(BytesIO(file.getvalue()))
                elif ext == 'pptx':
                    content = extract_text_from_pptx(BytesIO(file.getvalue()))
                elif ext == 'fcs':
                    content = parse_fcs_file(BytesIO(file.getvalue()))
                elif ext == 'hgb':
                    content = parse_hgb_file(BytesIO(file.getvalue()))
                elif ext in ['xlsx', 'xls']:
                    content = pd.read_excel(file)
                elif ext == 'csv':
                    content = pd.read_csv(file)
                elif ext in ['txt', 'json', 'md', 'py']:
                    content = file.getvalue().decode("utf-8")[:10000]
                elif ext in ['png', 'jpg', 'jpeg']:
                    img = Image.open(file)
                    content = f"Image: {img.size}px, {img.mode}"
                else:
                    content = "Unsupported file type"
                st.session_state.uploaded_files[file.name] = {"type": ext, "content": content}
            except Exception as e:
                st.error(f"❌ Error processing {file.name}: {str(e)}")

# ─── Visualizations ──────────────────────────────────────
def auto_plot_histograms(df, title="Histogram"):
    numeric_cols = df.select_dtypes(include='number').columns.tolist()
    if not numeric_cols:
        st.warning("No numeric columns to plot.")
        return
    st.subheader(f"📊 {title}")
    for col in numeric_cols:
        fig, ax = plt.subplots()
        ax.hist(df[col].dropna(), bins=30, color='skyblue', edgecolor='black')
        ax.set_title(f"Histogram: {col}")
        ax.set_xlabel(col)
        ax.set_ylabel("Frequency")
        st.pyplot(fig)

# ─── AI Analysis ─────────────────────────────────────────
def analyze_content(files_info, question=None):
    prompt = "You are an expert data analyst. The user has uploaded biomedical files. Provide a detailed analysis of each one. Specifically:\n"
    prompt += "- Determine what kind of biological data is present\n"
    prompt += "- Analyze column names and whether any are related to hemoglobin (e.g. HGB, HEMO, RBC, etc.)\n"
    prompt += "- Include detailed statistical summaries\n"
    prompt += "- Infer the likely purpose of the file (e.g. hemoglobin monitoring, blood count, flow cytometry)\n"
    prompt += "- If a question is asked, focus your answer on it\n\n"

    for fname, fdata in files_info.items():
        prompt += f"\n===== FILE: {fname} (type: {fdata['type']}) =====\n"

        if isinstance(fdata['content'], pd.DataFrame):
            df = fdata['content']
            prompt += "Columns:\n"
            prompt += ", ".join(df.columns) + "\n"
            prompt += "\nStats:\n"
            desc = df.describe().round(3).to_string()
            prompt += desc + "\n"

            # Detect hemoglobin-related terms
            hemo_terms = [col for col in df.columns if any(keyword in col.lower() for keyword in ['hgb', 'hemo', 'rbc'])]
            if hemo_terms:
                prompt += f"\n⚠️ Hemoglobin-related columns detected: {', '.join(hemo_terms)}\n"
            else:
                prompt += "\nNo hemoglobin-related column names detected.\n"

        elif isinstance(fdata['content'], str):
            sample_text = fdata['content'][:1000]
            prompt += f"File text sample:\n{sample_text}\n"

    if question:
        prompt += f"\nUser question: {question}\n"

    else:
        prompt += "\nGive your detailed interpretation of each file, what it contains, and whether it's likely to be related to hemoglobin analysis.\n"

    try:
        response = client.chat.completions.create(
            model=MODEL_OPTIONS[selected_model],
            messages=[{"role": "user", "content": prompt}],
            temperature=0.2,
            max_tokens=1500
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"AI Error: {str(e)}"

# ─── UI Layout ───────────────────────────────────────────
st.title("📁 File Analysis LLM Agent")

uploaded_files = st.file_uploader(
    "Upload files (FCS, HGB, Excel, PDF, etc.)",
    type=['fcs', 'hgb', 'pdf', 'txt', 'docx', 'pptx', 'csv', 'xlsx', 'xls', 'json', 'png', 'jpg', 'jpeg'],
    accept_multiple_files=True
)

if uploaded_files:
    process_uploaded_files(uploaded_files)

    with st.expander("📂 View Uploaded Files"):
        for name, data in st.session_state.uploaded_files.items():
            desc = "text" if isinstance(data['content'], str) else "dataframe"
            st.write(f"**{name}** ({data['type']}, {desc})")

    tab1, tab2, tab3 = st.tabs(["🔍 Summary", "❓ Ask Question", "📊 Visualize"])

    with tab1:
        if st.button("Generate Summary"):
            with st.spinner(f"Analyzing with {selected_model}..."):
                result = analyze_content(st.session_state.uploaded_files)
                st.session_state.history.append(("assistant", result))
                st.markdown(result)

    with tab2:
        q = st.text_input("Ask about your files:")
        if st.button("Get Answer") and q:
            with st.spinner(f"Researching with {selected_model}..."):
                answer = analyze_content(st.session_state.uploaded_files, q)
                st.session_state.history.append(("user", q))
                st.session_state.history.append(("assistant", answer))
                st.markdown(answer)

    with tab3:
        for fname, fdata in st.session_state.uploaded_files.items():
            if isinstance(fdata['content'], pd.DataFrame):
                st.markdown(f"### {fname}")
                auto_plot_histograms(fdata['content'])

st.divider()
st.subheader("💬 Conversation History")
for role, msg in st.session_state.history:
    st.markdown(f"**{role.capitalize()}**: {msg}")
if st.button("Clear History"):
    st.session_state.history = []
    st.rerun()
